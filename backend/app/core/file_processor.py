"""
RouteGen AI — File Processor

Extracts plain text from uploaded PDF, PPTX, and image files so it can be
chunked and indexed by the RAG store. Images are OCR'd with pytesseract when
the Tesseract binary is available; otherwise they're described by Gemini
Vision as a fallback.
"""

import os
import base64
import logging

logger = logging.getLogger("routegen.file_processor")

try:
    from pypdf import PdfReader
    HAS_PYPDF = True
except ImportError:
    HAS_PYPDF = False

try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

try:
    from PIL import Image
    import pytesseract
    HAS_OCR_LIBS = True
except ImportError:
    HAS_OCR_LIBS = False

IMAGE_TYPES = {"jpg", "jpeg", "png"}


def _extract_pdf(file_path: str) -> str:
    if not HAS_PYPDF:
        logger.error("pypdf not installed — cannot extract PDF text.")
        return ""
    try:
        reader = PdfReader(file_path)
        return "\n".join(page.extract_text() or "" for page in reader.pages)
    except Exception as e:
        logger.error(f"PDF extraction failed for {file_path}: {e}")
        return ""


def _extract_pptx(file_path: str) -> str:
    if not HAS_PPTX:
        logger.error("python-pptx not installed — cannot extract PPTX text.")
        return ""
    try:
        prs = Presentation(file_path)
        parts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame and shape.text_frame.text:
                    parts.append(shape.text_frame.text)
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text
                if notes:
                    parts.append(notes)
        return "\n".join(parts)
    except Exception as e:
        logger.error(f"PPTX extraction failed for {file_path}: {e}")
        return ""


def _extract_image_ocr(file_path: str) -> str:
    """Try local Tesseract OCR. Returns '' (not raises) if Tesseract isn't installed."""
    try:
        image = Image.open(file_path)
        return pytesseract.image_to_string(image)
    except Exception as e:
        logger.warning(f"Tesseract OCR unavailable/failed ({e}) — falling back to Gemini Vision.")
        return ""


async def _extract_image_gemini(file_path: str) -> str:
    """Fallback: describe the image via Gemini Vision when OCR is unavailable or empty."""
    try:
        import litellm

        with open(file_path, "rb") as f:
            image_b64 = base64.b64encode(f.read()).decode("utf-8")

        ext = os.path.splitext(file_path)[1].lstrip(".").lower()
        mime = f"image/{'jpeg' if ext in ('jpg', 'jpeg') else ext}"

        response = await litellm.acompletion(
            model="gemini/gemini-3.5-flash",
            messages=[{
                "role": "user",
                "content": [
                    {"type": "text", "text": "Describe everything visible in this image, including all text, in detail."},
                    {"type": "image_url", "image_url": {"url": f"data:{mime};base64,{image_b64}"}},
                ],
            }],
        )
        return response.choices[0].message.content or ""
    except Exception as e:
        logger.error(f"Gemini Vision extraction failed for {file_path}: {e}")
        return ""


async def extract_text(file_path: str, file_type: str) -> str:
    """
    Extract text from a file based on its type ("pdf", "pptx", "jpg"/"jpeg"/"png").
    Returns "" (with the error logged) if extraction fails or the type is unsupported.
    """
    file_type = file_type.lower().lstrip(".")
    try:
        if file_type == "pdf":
            return _extract_pdf(file_path)
        elif file_type == "pptx":
            return _extract_pptx(file_path)
        elif file_type in IMAGE_TYPES:
            text = _extract_image_ocr(file_path) if HAS_OCR_LIBS else ""
            if not text.strip():
                text = await _extract_image_gemini(file_path)
            return text
        else:
            logger.error(f"Unsupported file type for extraction: {file_type}")
            return ""
    except Exception as e:
        logger.error(f"extract_text failed for {file_path} ({file_type}): {e}")
        return ""
